import os
from pathlib import Path

from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

from backend.tools.base import Tool, ToolResult


BLOCKED_EXECUTABLE_EXTENSIONS = {}

TEXT_EXTENSIONS = {
    "txt", "py", "json", "md", "log", "js", "ts",
    "html", "css", "xml", "csv", "yml", "yaml", "ini", "cfg"
}


def get_extension(path: str) -> str:
    return path.lower().split(".")[-1]


def _resolve_path(path: str, must_exist: bool = False) -> tuple[Path | None, ToolResult | None]:
    try:
        if must_exist:
            resolved = Path(path).expanduser().resolve(strict=True)
        else:
            resolved = Path(path).expanduser().resolve()
    except OSError:
        return None, ToolResult(status="error", error="Invalid file path")

    return resolved, None


def _validate_openable_file(path: str) -> tuple[Path | None, ToolResult | None]:
    resolved, error = _resolve_path(path, must_exist=True)
    if error:
        return None, error

    if not resolved.is_file():
        return None, ToolResult(status="error", error="Path is not a file")

    if resolved.suffix.lower() in BLOCKED_EXECUTABLE_EXTENSIONS:
        return None, ToolResult(status="error", error="Opening this file type is blocked")

    return resolved, None


def read_text(path: str) -> ToolResult:
    try:
        with open(path, "r", encoding="utf-8") as file:
            return ToolResult(status="ok", content=file.read())
    except UnicodeDecodeError:
        return ToolResult(status="error", error="UTF-8 decode error")


def read_docx(path: str) -> ToolResult:
    doc = Document(path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return ToolResult(status="ok", content=text)


def read_pptx(path: str) -> ToolResult:
    prs = Presentation(path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)

    return ToolResult(status="ok", content="\n".join(text))


def read_xlsx(path: str) -> ToolResult:
    wb = load_workbook(path, data_only=True)
    result = {}

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows = []

        for row in sheet.iter_rows(values_only=True):
            rows.append(["" if cell is None else str(cell) for cell in row])

        result[sheet_name] = rows

    return ToolResult(status="ok", content=result)


def read_file(path: str) -> ToolResult:
    if not os.path.exists(path):
        return ToolResult(status="error", error="File path does not exist")

    ext = get_extension(path)

    if ext in TEXT_EXTENSIONS:
        return read_text(path)

    if ext == "docx":
        return read_docx(path)

    if ext == "pptx":
        return read_pptx(path)

    if ext == "xlsx":
        return read_xlsx(path)

    return ToolResult(status="error", error=f"Unsupported file type: .{ext}")


def list_files(path: str) -> ToolResult:
    if not os.path.exists(path):
        return ToolResult(status="error", error="Directory does not exist")

    if not os.path.isdir(path):
        return ToolResult(status="error", error="Path is not a directory")

    return ToolResult(status="ok", content=os.listdir(path))


def open_file(path: str) -> ToolResult:
    resolved, error = _validate_openable_file(path)
    if error:
        return error

    os.startfile(str(resolved))
    return ToolResult(status="ok", content="Opened file")


def create_file(path: str, content: str = "") -> ToolResult:
    resolved, error = _resolve_path(path)
    if error:
        return error

    if resolved.exists():
        return ToolResult(status="error", error="File already exists")

    resolved.parent.mkdir(parents=True, exist_ok=True)
    resolved.write_text(content, encoding="utf-8")
    return ToolResult(status="ok", content="Created file")


def write_file(path: str, content: str) -> ToolResult:
    resolved, error = _resolve_path(path)
    if error:
        return error

    if resolved.exists() and not resolved.is_file():
        return ToolResult(status="error", error="Path is not a file")

    resolved.parent.mkdir(parents=True, exist_ok=True)
    resolved.write_text(content, encoding="utf-8")
    return ToolResult(status="ok", content="Wrote file")


def delete_file(path: str) -> ToolResult:
    resolved, error = _resolve_path(path, must_exist=True)
    if error:
        return error

    if not resolved.is_file():
        return ToolResult(status="error", error="Path is not a file")

    resolved.unlink()
    return ToolResult(status="ok", content="Deleted file")


FILE_TOOLS = [
    Tool(
        name="read_file",
        description="Read the content of a local file.",
        parameters={
            "type": "object",
            "properties": {
                "path": {"type": "string"}
            },
            "required": ["path"],
        },
        function=read_file,
    ),
    Tool(
        name="list_files",
        description="List files and folders inside a local directory.",
        parameters={
            "type": "object",
            "properties": {
                "path": {"type": "string"}
            },
            "required": ["path"],
        },
        function=list_files,
    ),
    Tool(
        name="open_file",
        description="Open a local file.",
        parameters={
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"],
        },
        function=open_file,
    ),
    Tool(
        name="create_file",
        description="Create a new UTF-8 file.",
        parameters={
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "content": {"type": "string"},
            },
            "required": ["path"],
        },
        function=create_file,
    ),
    Tool(
        name="write_file",
        description="Write a UTF-8 file.",
        parameters={
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "content": {"type": "string"},
            },
            "required": ["path", "content"],
        },
        function=write_file,
    ),
    Tool(
        name="delete_file",
        description="Delete a file after user confirmation.",
        parameters={
            "type": "object",
            "properties": {
                "path": {"type": "string"}
            },
            "required": ["path"],
        },
        function=delete_file,
    ),
]
