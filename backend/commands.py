import os
import ctypes
import webbrowser

from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

def error(message: str):
    return {
        "status": "error",
        "content": message
    }



# Read/Write files ----------------------------------------
def read_file(path: str):
    try:
        if not os.path.exists(path):
            return error(f"File path doesn't exist")

        ext = get_extension(path)

        result = read_dispatch(path, ext)
        return result

    except Exception as e:
        return error(str(e))

def get_extension(path: str) -> str:
    return path.lower().split(".")[-1]


def read_dispatch(path: str, ext: str):
    if ext in ["txt", "py", "json", "md", "log", "js", "ts", "html", "css", "xml", "csv", "yml", "yaml", "ini", "cfg"]:
        return read_text(path)
    elif ext == "docx":
        return read_docx(path)
    elif ext == "pptx":
        return read_pptx(path)
    elif ext == "xlsx":
        return read_xlsx(path)
    else:
        return error(f"Unsupported file type: .{ext}")


def read_text(path: str):
    try:
        with open(path, "r", encoding="utf-8") as f:
            return {
                "status": "ok",
                "type": "text",
                "content": f.read()
            }
    except UnicodeDecodeError:
        return error("UTF-8 decode error")
    except Exception as e:
        return error(f"Text read error: {str(e)}")

def read_docx(path: str):
    try:
        doc = Document(path)
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        return {
            "status": "ok",
            "type": "docx",
            "content": text
        }

    except Exception as e:
        return error(f"DOCX read error: {str(e)}")

def read_pptx(path: str):
    try:
        prs = Presentation(path)
        text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)

        return {
            "status": "ok",
            "type": "pptx",
            "content": "\n".join(text)
        }
    except Exception as e:
        return error(f"PPTX read error: {str(e)}")
def read_xlsx(path: str):
    try:
        wb = load_workbook(path, data_only=True)

        result = {
            "status": "ok",
            "type": "xlsx",
            "sheets": {}
        }

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []

            for row in sheet.iter_rows(values_only=True):
                clean_row = [
                    "" if cell is None else str(cell)
                    for cell in row
                ]
                rows.append(clean_row)

            result["sheets"][sheet_name] = rows

        return result

    except Exception as e:
        return error(f"XLSX read error: {str(e)}")

# List files ----------------------------------------
def list_files(path):
    try:
        return {"status": "ok", "content": os.listdir(path)}
    except Exception as e:
        return error(str(e))


# Open URL ----------------------------------------
def open_url(url):
    webbrowser.open(url)

# Open file ----------------------------------------
def open_file(path):
    try:
        if not os.path.exists(path):
            return error(f"File path doesn't exist")
        os.startfile(path)
        return {"status": "ok"}
    except Exception as e:
        return error(str(e))

# Set volume ----------------------------------------
def set_volume(direction, amount):
    if direction == "up":
        for i in range(amount):
            pressKey(0xAF)
    elif direction == "down":
        for i in range(amount):
            pressKey(0xAE)

# Press key ----------------------------------------
def pressKey(key):
    ctypes.windll.user32.keybd_event(key, 0, 0, 0)
    ctypes.windll.user32.keybd_event(key, 0, 2, 0)